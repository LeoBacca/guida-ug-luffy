# -*- coding: utf-8 -*-
"""
UG (Blue/Green) LUFFY — OP16 PRISON BREAK · TACTICAL OPS MANUAL
Guida PPTX rifatta con design system "Impel Down — Tactical HUD"
(blueprint carbone + colore-codice blu/verde, oro=Hancock, immagini carte reali).

Genera: "Guida UG Luffy - OP16 (TACTICAL).pptx"
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.oxml.ns import qn
from PIL import Image, ImageDraw
import assets as A

BASE = os.path.dirname(os.path.abspath(__file__))
IMG  = os.path.join(BASE, "assets", "IMG")

# ====================== PALETTE (Tactical HUD) ======================
BG_CARBON  = RGBColor(0x0A,0x10,0x18)
BG_DEEP    = RGBColor(0x06,0x08,0x0D)
SURFACE    = RGBColor(0x10,0x1A,0x28)
SURFACE_HI = RGBColor(0x16,0x24,0x3A)
GRID_MAJOR = RGBColor(0x1B,0x2C,0x42)
GRID_MINOR = RGBColor(0x13,0x20,0x2F)
INK        = RGBColor(0xEA,0xF2,0xFF)
INK_DIM    = RGBColor(0x9D,0xB2,0xCC)
MUTED      = RGBColor(0x5A,0x6E,0x88)
BLUE       = RGBColor(0x2E,0x8B,0xFF)   # VALORE / pescata / flusso
BLUE_DEEP  = RGBColor(0x0E,0x5B,0xD6)
GREEN      = RGBColor(0x27,0xE0,0x8A)   # TEMPO / ramp
GREEN_DEEP = RGBColor(0x0F,0x9D,0x5C)
GOLD       = RGBColor(0xFF,0xC2,0x3D)   # SOLO Boa Hancock / win condition
DANGER     = RGBColor(0xFF,0x54,0x66)   # rischio / counterplay / hard
CYAN       = RGBColor(0x3F,0xE0,0xFF)   # hairline / corner-bracket / crosshair
GRAD_MID   = RGBColor(0x14,0x30,0x6B)

MONO = "Consolas"
DISP = "Arial Black"
UI   = "Segoe UI"
UISB = "Segoe UI Semibold"

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]
PNUM = [0]


# ====================== LOW-LEVEL HELPERS ======================
def _hex(c): return "%02X%02X%02X" % (c[0], c[1], c[2])


def fill_alpha(shp, color, alpha):
    """Fill solido con alpha (0-100). alpha=100 opaco."""
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    sp = shp.fill._xPr.find(qn('a:solidFill'))
    srgb = sp.find(qn('a:srgbClr'))
    a = srgb.makeelement(qn('a:alpha'), {'val': str(int(alpha*1000))})
    srgb.append(a)


def _no_line(shp): shp.line.fill.background()


def _line(shp, color, w):
    shp.line.color.rgb = color
    shp.line.width = Pt(w)


def _shadow(shp, blur=70000, dist=38000, alpha=46000, dirn=5400000):
    el = shp._element.spPr
    ef = el.makeelement(qn('a:effectLst'), {})
    sh = ef.makeelement(qn('a:outerShdw'),
                        {'blurRad': str(blur), 'dist': str(dist),
                         'dir': str(dirn), 'rotWithShape': '0'})
    clr = sh.makeelement(qn('a:srgbClr'), {'val': '02060C'})
    al = clr.makeelement(qn('a:alpha'), {'val': str(alpha)})
    clr.append(al); sh.append(clr); ef.append(sh); el.append(ef)


def rect(s, x, y, w, h, fill=None, line=None, lw=0.75, alpha=None,
         shape=MSO_SHAPE.RECTANGLE, shadow=False, radius=None):
    shp = s.shapes.add_shape(shape, int(x), int(y), int(w), int(h))
    shp.shadow.inherit = False
    if fill is None:
        shp.fill.background()
    elif alpha is not None:
        fill_alpha(shp, fill, alpha)
    else:
        shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None:
        _no_line(shp)
    else:
        _line(shp, line, lw)
    if shadow:
        _shadow(shp)
    if radius is not None and shape == MSO_SHAPE.ROUNDED_RECTANGLE:
        try: shp.adjustments[0] = radius
        except Exception: pass
    return shp


def _track(run, spc):
    """Tracking (lettere spaziate) via a:spc in milliemi (es +220 -> 2200)."""
    rPr = run._r.get_or_add_rPr()
    rPr.set('spc', str(int(spc)))


def _add_runs(p, segs, dcolor, dsz, dbold, dfont):
    for seg in segs:
        txt, ov = (seg if isinstance(seg, tuple) else (seg, {}))
        r = p.add_run(); r.text = txt
        f = r.font
        f.size = Pt(ov.get("sz", dsz))
        f.bold = ov.get("bold", dbold)
        f.italic = ov.get("italic", False)
        f.name = ov.get("font", dfont)
        f.color.rgb = ov.get("color", dcolor)
        if ov.get("track"):
            _track(r, ov["track"])


def text(s, x, y, w, h, lines, color=INK, sz=13, bold=False, font=UI,
         align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, space_after=4,
         line_spacing=1.12, wrap=True, track=None):
    tb = s.shapes.add_textbox(int(x), int(y), int(w), int(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    for m in ("left", "right", "top", "bottom"):
        setattr(tf, f"margin_{m}", Pt(3))
    first = True
    for ln in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align
        p.space_after = Pt(space_after)
        p.space_before = Pt(0)
        p.line_spacing = line_spacing
        if isinstance(ln, dict):
            segs = ln.get("segs", "")
            if isinstance(segs, str): segs = [segs]
            p.level = ln.get("level", 0)
            if ln.get("align"): p.alignment = ln["align"]
            if ln.get("space_after") is not None:
                p.space_after = Pt(ln["space_after"])
            if ln.get("line_spacing"):
                p.line_spacing = ln["line_spacing"]
            tr = ln.get("track", track)
            base = [{**({} if not isinstance(sg, tuple) else sg[1])} for sg in segs]
            # applica track di default a tutti i run del paragrafo se richiesto
            segs2 = []
            for sg in segs:
                t2, o2 = (sg if isinstance(sg, tuple) else (sg, {}))
                if tr and "track" not in o2: o2 = {**o2, "track": tr}
                segs2.append((t2, o2))
            _add_runs(p, segs2, ln.get("color", color), ln.get("sz", sz),
                      ln.get("bold", bold), ln.get("font", font))
        else:
            segs = ln if isinstance(ln, list) else [ln]
            segs2 = []
            for sg in segs:
                t2, o2 = (sg if isinstance(sg, tuple) else (sg, {}))
                if track and "track" not in o2: o2 = {**o2, "track": track}
                segs2.append((t2, o2))
            _add_runs(p, segs2, color, sz, bold, font)
    return tb


def shape_text(shp, lines, color=INK, sz=13, bold=False, font=UI,
               align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, space_after=2,
               line_spacing=1.04, track=None):
    tf = shp.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for m in ("left", "right", "top", "bottom"):
        setattr(tf, f"margin_{m}", Pt(3))
    if isinstance(lines, str): lines = [lines]
    first = True
    for ln in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align
        p.space_after = Pt(space_after)
        p.line_spacing = line_spacing
        if isinstance(ln, dict):
            segs = ln.get("segs", "")
            if isinstance(segs, str): segs = [segs]
            if ln.get("align"): p.alignment = ln["align"]
            tr = ln.get("track", track)
            segs2 = []
            for sg in segs:
                t2, o2 = (sg if isinstance(sg, tuple) else (sg, {}))
                if tr and "track" not in o2: o2 = {**o2, "track": tr}
                segs2.append((t2, o2))
            _add_runs(p, segs2, ln.get("color", color), ln.get("sz", sz),
                      ln.get("bold", bold), ln.get("font", font))
        else:
            segs = ln if isinstance(ln, list) else [ln]
            segs2 = []
            for sg in segs:
                t2, o2 = (sg if isinstance(sg, tuple) else (sg, {}))
                if track and "track" not in o2: o2 = {**o2, "track": track}
                segs2.append((t2, o2))
            _add_runs(p, segs2, color, sz, bold, font)
    return shp


# ====================== BLUEPRINT GRID PNG ======================
GRID_PNG = os.path.join(IMG, "grid_blueprint.png")
def grid_png():
    if os.path.exists(GRID_PNG):
        return GRID_PNG
    W, H = 2667, 1500  # 200 dpi @ 13.333x7.5
    im = Image.new("RGB", (W, H), tuple(BG_CARBON))
    d = ImageDraw.Draw(im)
    minor = int(W/40)   # ~0.33in
    for x in range(0, W, minor):
        d.line([(x,0),(x,H)], fill=tuple(GRID_MINOR), width=1)
    for y in range(0, H, minor):
        d.line([(0,y),(W,y)], fill=tuple(GRID_MINOR), width=1)
    major = int(W/13)   # ~1in
    for x in range(0, W, major):
        d.line([(x,0),(x,H)], fill=tuple(GRID_MAJOR), width=1)
    for y in range(0, H, major):
        d.line([(0,y),(W,y)], fill=tuple(GRID_MAJOR), width=1)
    im.save(GRID_PNG)
    return GRID_PNG


def slide():
    return prs.slides.add_slide(BLANK)


def base(s, vignette=True):
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = BG_CARBON
    s.shapes.add_picture(grid_png(), 0, 0, SW, SH)
    if vignette:
        # vignette: 4 fasce scure ai bordi
        rect(s, 0, 0, SW, Inches(0.5), BG_DEEP, alpha=55)
        rect(s, 0, SH-Inches(0.5), SW, Inches(0.5), BG_DEEP, alpha=70)
        rect(s, 0, 0, Inches(0.35), SH, BG_DEEP, alpha=45)
        rect(s, SW-Inches(0.35), 0, Inches(0.35), SH, BG_DEEP, alpha=45)


# ====================== IMAGE / FRAME HELPERS ======================
def _imsize(path):
    try: return Image.open(path).size
    except Exception: return (1, 1)


def pic_contain(s, path, x, y, w, h):
    iw, ih = _imsize(path)
    ar = iw/ih; boxar = w/h
    if ar > boxar: nw, nh = w, int(w/ar)
    else: nh, nw = h, int(h*ar)
    nx = int(x + (w-nw)/2); ny = int(y + (h-nh)/2)
    s.shapes.add_picture(path, nx, ny, int(nw), int(nh))
    return nx, ny, int(nw), int(nh)


def pic_cover(s, path, x, y, w, h):
    """COVER crop: riempie il box, eccesso fuori (clip via crop nativo)."""
    iw, ih = _imsize(path)
    ar = iw/ih; boxar = w/h
    pic = s.shapes.add_picture(path, int(x), int(y), int(w), int(h))
    if ar > boxar:
        crop = (1 - boxar/ar)/2
        pic.crop_left = crop; pic.crop_right = crop
    else:
        crop = (1 - ar/boxar)/2
        pic.crop_top = crop; pic.crop_bottom = crop
    return pic


def corner_brackets(s, x, y, w, h, color=CYAN, seg=None, lw=1.6, gold=False):
    if gold: color = GOLD
    if seg is None: seg = Inches(0.22)
    t = Pt(lw)
    def Lc(cx, cy, dx, dy):
        # segmento orizzontale + verticale che formano una L all'angolo
        rect(s, cx if dx > 0 else cx-seg, cy if dy > 0 else cy-t, seg, t, color)
        rect(s, cx if dx > 0 else cx-t, cy if dy > 0 else cy-seg, t, seg, color)
    Lc(x,     y,     1,  1)
    Lc(x+w,   y,    -1,  1)
    Lc(x,     y+h,   1, -1)
    Lc(x+w,   y+h,  -1, -1)


def crosshair(s, cx, cy, size=None, color=CYAN):
    if size is None: size = Inches(0.16)
    rect(s, cx-size/2, cy-Pt(0.6), size, Pt(1.2), color)
    rect(s, cx-Pt(0.6), cy-size/2, Pt(1.2), size, color)


def scanline(s, x, y, w, code, name=None, color=CYAN):
    h = Inches(0.26)
    bar = rect(s, x, y, w, h, SURFACE, alpha=88)
    rect(s, x, y, Inches(0.05), h, color)
    segs = [(f"  {code}", {"font": MONO, "sz": 9.5, "color": color, "bold": True})]
    if name:
        segs.append((f"  ·  {name}", {"font": MONO, "sz": 9, "color": INK_DIM}))
    text(s, x+Inches(0.05), y, w-Inches(0.1), h, [segs],
         anchor=MSO_ANCHOR.MIDDLE, space_after=0)


def scrim_band(s, x, y, w, h, direction="up"):
    """Stack di rettangoli BG_DEEP con alpha decrescente per leggibilita testo."""
    n = 4
    alphas = [78, 55, 32, 0] if direction == "up" else [0, 32, 55, 78]
    for i in range(n):
        if alphas[i] == 0: continue
        rh = h/n
        ry = y + (h - rh*(i+1)) if direction == "up" else y + rh*i
        rect(s, x, ry, w, rh+Pt(1), BG_DEEP, alpha=alphas[i])


def art_band(s, key, x, y, w, h, code=None, lead=False, scrim=True):
    path = A.lead_art(key) if lead else A.card_art(key)
    pic_cover(s, path, x, y, w, h)
    if scrim:
        scrim_band(s, x, y+h*0.45, w, h*0.55, "up")
    if code:
        scanline(s, x+Inches(0.12), y+h-Inches(0.34), Inches(3.0), code)


def card_spotlight(s, key, x, y, w, h, gold=False, code=None):
    path = A.card_full(key)
    if gold:  # alone faked dietro
        rect(s, x-Inches(0.12), y-Inches(0.12), w+Inches(0.24), h+Inches(0.24),
             GOLD, alpha=16, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.05)
    nx, ny, nw, nh = pic_contain(s, path, x, y, w, h)
    corner_brackets(s, nx-Pt(3), ny-Pt(3), nw+Pt(6), nh+Pt(6),
                    gold=gold, seg=Inches(0.26), lw=2.0)
    if code:
        scanline(s, nx, ny+nh-Inches(0.30), nw, code, color=(GOLD if gold else CYAN))
    return nx, ny, nw, nh


# ====================== HUD MODULES ======================
def module(s, x, y, w, h, title=None, accent=CYAN, hi=False, brackets=False,
           tcolor=None):
    fill = SURFACE_HI if hi else SURFACE
    m = rect(s, x, y, w, h, fill, line=GRID_MAJOR, lw=0.75,
             shape=MSO_SHAPE.RECTANGLE, shadow=True)
    rect(s, x, y, Inches(0.055), h, accent)   # rail accent sx
    hh = Inches(0)
    if title:
        hh = Inches(0.42)
        text(s, x+Inches(0.2), y+Inches(0.06), w-Inches(0.3), hh,
             [{"segs": title, "track": 130}], color=tcolor or accent, sz=11.5,
             bold=True, font=UISB, anchor=MSO_ANCHOR.MIDDLE)
        rect(s, x+Inches(0.2), y+hh, w-Inches(0.4), Pt(0.8), GRID_MAJOR)
    if brackets:
        corner_brackets(s, x, y, w, h, accent, seg=Inches(0.18), lw=1.3)
    return m, y+hh


def chip(s, x, y, w, h, label, fill=SURFACE, txt=INK, line=None, sz=11,
         bold=True, font=UISB, track=80):
    b = rect(s, x, y, w, h, fill, line=line, lw=1.0,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.5)
    shape_text(b, label, color=txt, sz=sz, bold=bold, font=font, track=track)
    return b


def attr_chip(s, x, y, w, kind):
    col = BLUE if kind == "BLUE" else GREEN
    h = Inches(0.34)
    b = rect(s, x, y, w, h, SURFACE, line=col, lw=1.25,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.5)
    rect(s, x+Inches(0.1), y+h/2-Inches(0.07), Inches(0.14), Inches(0.14), col,
         shape=MSO_SHAPE.OVAL)
    shape_text(b, [{"segs": kind, "align": PP_ALIGN.LEFT}], color=col, sz=10.5,
               bold=True, font=MONO, align=PP_ALIGN.LEFT, track=60)
    # sposta testo a dx del dot
    return b


def cost_chip(s, x, y, key, extra=None):
    c = A.cost(key)
    label = f"COST // {c}" + (f"  {extra}" if extra else "")
    w = Inches(1.5 + (0.5 if extra else 0))
    col = GOLD if key == "hancock" else BLUE
    b = rect(s, x, y, w, Inches(0.34), SURFACE, line=col, lw=1.1,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.4)
    shape_text(b, label, color=col, sz=10, bold=True, font=MONO, track=40)
    return b


def led_bar(s, x, y, w, h, filled, total, color=BLUE, peak_gold=False,
            danger_thr=None):
    gap = Pt(2)
    seg_w = (w - gap*(total-1)) / total
    for i in range(total):
        sx = x + (seg_w+gap)*i
        on = i < filled
        c = color
        if on and peak_gold and i == filled-1: c = GOLD
        if on and danger_thr is not None and i < danger_thr: c = DANGER
        rect(s, sx, y, seg_w, h, c if on else BG_DEEP,
             line=(None if on else GRID_MAJOR), lw=0.5)


def led_dot(s, x, y, state="blue", size=None):
    if size is None: size = Inches(0.13)
    col = {"green": GREEN, "blue": BLUE, "gold": GOLD, "danger": DANGER,
           "muted": MUTED}.get(state, BLUE)
    rect(s, x, y, size, size, col, shape=MSO_SHAPE.OVAL)


def node_square(s, cx, cy, size=Pt(7), color=CYAN):
    rect(s, cx-size/2, cy-size/2, size, size, color)


def arrow_flow(s, x1, y1, x2, y2, color=BLUE, w=2.4, elbow="h"):
    """Connettore ortogonale a gomito con chevron + nodi quadrati."""
    if elbow == "h":  # H poi V
        rect(s, min(x1,x2), y1-Pt(w/2), abs(x2-x1), Pt(w), color)
        rect(s, x2-Pt(w/2), min(y1,y2), Pt(w), abs(y2-y1), color)
    else:             # V poi H
        rect(s, x1-Pt(w/2), min(y1,y2), Pt(w), abs(y2-y1), color)
        rect(s, min(x1,x2), y2-Pt(w/2), abs(x2-x1), Pt(w), color)
    # chevron a destinazione
    chev(s, x2, y2, color, direction=("down" if y2>y1 and elbow=="h" else
                                       "right"))
    node_square(s, x1, y1, color=color)


def chev(s, x, y, color=BLUE, direction="right", size=Inches(0.13)):
    sh = MSO_SHAPE.ISOSCELES_TRIANGLE
    t = rect(s, x-size/2, y-size/2, size, size, color, shape=sh)
    rot = {"right": 90, "left": 270, "down": 180, "up": 0}[direction]
    t.rotation = rot
    return t


def straight_arrow(s, x1, y1, x2, y2, color=BLUE, w=2.4):
    c = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, int(x1),int(y1),int(x2),int(y2))
    c.line.color.rgb = color; c.line.width = Pt(w); c.shadow.inherit = False
    ln = c.line._get_or_add_ln()
    ln.append(ln.makeelement(qn('a:tailEnd'),
              {'type':'triangle','w':'med','len':'med'}))
    return c


# ====================== DISPLAY / NUMERALS ======================
def giant(s, x, y, txt, sz=300, color=GRID_MAJOR, align=PP_ALIGN.LEFT):
    # larghezza proporzionale al testo (Arial Black ~0.62em/char), clamp al bordo
    bw = Emu(int(sz * 0.62 * max(1, len(txt)) * 12700))
    if align == PP_ALIGN.LEFT and x + bw > SW - Inches(0.1):
        bw = max(Inches(1), SW - Inches(0.1) - x)
    text(s, x, y, bw, Inches(sz/55), [txt], color=color, sz=sz,
         bold=True, font=DISP, align=align, anchor=MSO_ANCHOR.MIDDLE,
         line_spacing=0.9, wrap=False)


def data_big(s, x, y, w, value, label, color=BLUE, sub=None, vsz=60):
    text(s, x, y, w, Inches(0.26), [{"segs": label, "track": 200}], color=MUTED,
         sz=9, bold=True, font=MONO, align=PP_ALIGN.CENTER)
    text(s, x, y+Inches(0.22), w, Inches(vsz/55), [value], color=color, sz=vsz,
         bold=True, font=DISP, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
         line_spacing=0.9)
    if sub:
        text(s, x, y+Inches(0.22+vsz/72), w, Inches(0.26),
             [{"segs": sub, "track": 120}], color=INK_DIM, sz=9, font=MONO,
             align=PP_ALIGN.CENTER)


def gradient_band(s, y, h=Inches(0.12)):
    n = 24
    cw = SW / n
    for i in range(n):
        t = i/(n-1)
        if t < 0.5:
            tt = t/0.5
            c = RGBColor(int(BLUE_DEEP[0]+(GRAD_MID[0]-BLUE_DEEP[0])*tt),
                         int(BLUE_DEEP[1]+(GRAD_MID[1]-BLUE_DEEP[1])*tt),
                         int(BLUE_DEEP[2]+(GRAD_MID[2]-BLUE_DEEP[2])*tt))
        else:
            tt = (t-0.5)/0.5
            c = RGBColor(int(GRAD_MID[0]+(GREEN_DEEP[0]-GRAD_MID[0])*tt),
                         int(GRAD_MID[1]+(GREEN_DEEP[1]-GRAD_MID[1])*tt),
                         int(GRAD_MID[2]+(GREEN_DEEP[2]-GRAD_MID[2])*tt))
        rect(s, cw*i, y, cw+Emu(9525), h, c)


def term_log(s, x, y, w, lines, prefix="> ", title=None, accent=GREEN):
    yy = y
    if title:
        text(s, x, yy, w, Inches(0.26), [{"segs": title, "track": 160}],
             color=accent, sz=10, bold=True, font=MONO)
        yy += Inches(0.3)
    segs = []
    for ln in lines:
        if isinstance(ln, tuple):
            t2, col = ln
        else:
            t2, col = ln, INK_DIM
        segs.append({"segs": [(prefix, {"color": accent, "font": MONO}),
                              (t2, {"color": col, "font": MONO})],
                     "sz": 10.5, "space_after": 3})
    text(s, x, yy, w, Inches(0.32)*len(lines)+Inches(0.1), segs, line_spacing=1.15)


# ====================== HEADER / FOOTER ======================
CHAP_ACCENT = {1:BLUE, 2:GREEN, 3:BLUE, 4:GREEN, 5:BLUE, 6:GREEN}

def kicker_h2(s, kicker, title, folio=None, accent=CYAN, sub=None):
    text(s, Inches(0.55), Inches(0.42), Inches(10), Inches(0.3),
         [{"segs": kicker, "track": 220}], color=accent, sz=10.5, bold=True,
         font=UISB)
    text(s, Inches(0.52), Inches(0.66), Inches(11.4), Inches(0.62),
         [title], color=INK, sz=26, bold=True, font=UISB)
    if sub:
        text(s, Inches(0.55), Inches(1.26), Inches(11.6), Inches(0.4),
             [sub], color=INK_DIM, sz=13.5, font=UI)
    rect(s, Inches(0.55), Inches(1.2) if not sub else Inches(1.66),
         Inches(12.25), Pt(1), GRID_MAJOR)
    rect(s, Inches(0.55), Inches(1.2) if not sub else Inches(1.66),
         Inches(2.2), Pt(1.4), accent)
    if folio:
        text(s, Inches(10.0), Inches(0.46), Inches(2.78), Inches(0.4),
             [{"segs": folio, "track": 160, "align": PP_ALIGN.RIGHT}],
             color=MUTED, sz=9.5, bold=True, font=MONO, align=PP_ALIGN.RIGHT)


def footer_hud(s, chap=None, coord="OP16 // PRISON BREAK"):
    PNUM[0] += 1
    n = PNUM[0]
    gradient_band(s, SH-Inches(0.06), Inches(0.06))
    text(s, Inches(0.45), Inches(7.04), Inches(2.0), Inches(0.32),
         [{"segs": f"P. {n:02d}", "track": 140}], color=MUTED, sz=9, bold=True,
         font=MONO)
    text(s, Inches(4.5), Inches(7.04), Inches(4.3), Inches(0.32),
         [{"segs": coord, "track": 180, "align": PP_ALIGN.CENTER}], color=MUTED,
         sz=9, font=MONO, align=PP_ALIGN.CENTER)
    if chap:
        for i in range(1, 7):
            cx = Inches(11.3) + Inches(0.27)*(i-1)
            on = (i == chap)
            col = CHAP_ACCENT[i] if on else GRID_MAJOR
            rect(s, cx, Inches(7.08), Inches(0.2), Inches(0.2),
                 col if on else None, line=(None if on else MUTED), lw=0.75)
    return n


# ============================================================
#  SLIDE 01 — COVER / BOOT SCREEN
# ============================================================
def cover():
    s = slide(); base(s, vignette=False)
    # art band hero (leader luffy) full-width centrale
    art_band(s, "luffy", 0, Inches(1.35), SW, Inches(3.7), lead=True, scrim=True)
    rect(s, 0, Inches(1.35), SW, Inches(3.7), BG_DEEP, alpha=18)
    # top HUD strip
    rect(s, 0, 0, SW, Inches(0.5), BG_DEEP, alpha=80)
    text(s, Inches(0.55), Inches(0.12), Inches(8), Inches(0.3),
         [{"segs": "ONE PIECE CARD GAME // OP16 PRISON BREAK // TACTICAL DECK MANUAL",
           "track": 200}], color=CYAN, sz=10, bold=True, font=MONO)
    text(s, Inches(10.4), Inches(0.12), Inches(2.4), Inches(0.3),
         [{"segs": "REV 06.2026", "track": 160, "align": PP_ALIGN.RIGHT}],
         color=MUTED, sz=9.5, font=MONO, align=PP_ALIGN.RIGHT)
    # masthead
    text(s, Inches(0.7), Inches(1.55), Inches(12), Inches(1.7),
         [{"segs": [("UG ", {"color": BLUE}), ("LUFFY", {"color": GREEN})],
           "track": -10}], sz=92, bold=True, font=DISP, anchor=MSO_ANCHOR.MIDDLE,
         wrap=False)
    chip(s, Inches(0.78), Inches(3.15), Inches(5.7), Inches(0.5),
         "BLUE / GREEN  ·  IMPEL DOWN LUFFY", SURFACE, INK, line=CYAN, sz=13,
         track=120)
    # attr chips
    attr_chip(s, Inches(0.78), Inches(3.78), Inches(1.5), "BLUE")
    attr_chip(s, Inches(2.4),  Inches(3.78), Inches(1.5), "GREEN")
    chip(s, Inches(4.02), Inches(3.78), Inches(2.45), Inches(0.34),
         "LEADER // OP16-022", SURFACE, BLUE, line=BLUE, sz=10, track=40, font=MONO)
    # lower module: tagline + author
    rect(s, 0, Inches(5.35), SW, Inches(2.15), BG_DEEP, alpha=82)
    rect(s, 0, Inches(5.35), SW, Pt(1.2), CYAN)
    text(s, Inches(0.7), Inches(5.55), Inches(8.4), Inches(0.9),
         [{"segs": "« THE AGGRO KING IS BACK »", "track": 40}], color=GOLD, sz=24,
         bold=True, font=UISB)
    text(s, Inches(0.72), Inches(6.12), Inches(8.6), Inches(0.5),
         ["Aggro a valore sostenuto — guida dettagliata, carta per carta, matchup per matchup."],
         color=INK_DIM, sz=14)
    # author card right
    m,_ = module(s, Inches(9.3), Inches(5.55), Inches(3.45), Inches(1.5),
                 title="AUTHOR // INTEL")
    text(s, Inches(9.5), Inches(6.0), Inches(3.1), Inches(1.0),
         [{"segs": "Enricomaria Rustico", "sz": 14, "bold": True, "color": INK,
           "space_after": 2},
          {"segs": "Dogofeisdom · @DogXWisdom", "sz": 11, "color": BLUE,
           "font": MONO, "space_after": 4},
          {"segs": "Bandai Card Fest 2024 · Top 8 World 2025", "sz": 10,
           "color": MUTED}])
    gradient_band(s, SH-Inches(0.14), Inches(0.14))
    PNUM[0] += 1


# ============================================================
#  SLIDE 02 — INDICE
# ============================================================
def indice():
    s = slide(); base(s)
    kicker_h2(s, "TABLE OF CONTENTS // 6 SECTORS", "INDICE OPERATIVO",
              folio="NAV 00/06")
    giant(s, Inches(8.6), Inches(3.6), "06", sz=300, color=GRID_MAJOR,
          align=PP_ALIGN.LEFT)
    items = [
        ("01","THE AGGRO KING IS BACK","Identità del mazzo · stile di gioco",1),
        ("02","META ANALYSIS","Meta OP16 · tier list · matrice win-rate",2),
        ("03","THE DECKLIST","Lista · carta per carta · tech pool",3),
        ("04","GENERAL STRATEGY","Mulligan · early / mid / late · combo-line",4),
        ("05","THE MATCHUPS BIBLE","6 matchup chiave + life management",5),
        ("06","CONCLUSION","Chiusura · nuova lista in test",6),
    ]
    y0 = Inches(1.95); rh = Inches(0.78); gap = Inches(0.08)
    for i,(n,t,sub,ch) in enumerate(items):
        yy = y0 + (rh+gap)*i
        acc = CHAP_ACCENT[ch]
        m = rect(s, Inches(0.55), yy, Inches(8.0), rh, SURFACE, line=GRID_MAJOR,
                 lw=0.75, shadow=True)
        rect(s, Inches(0.55), yy, Inches(0.06), rh, acc)
        text(s, Inches(0.78), yy, Inches(1.1), rh, [n], color=acc, sz=30,
             bold=True, font=DISP, anchor=MSO_ANCHOR.MIDDLE)
        led_dot(s, Inches(1.95), yy+rh/2-Inches(0.065), state=("blue" if acc==BLUE else "green"))
        text(s, Inches(2.25), yy+Inches(0.1), Inches(6.1), Inches(0.36),
             [{"segs": t, "track": 40}], color=INK, sz=16, bold=True, font=UISB)
        text(s, Inches(2.25), yy+Inches(0.44), Inches(6.1), Inches(0.3),
             [sub], color=INK_DIM, sz=11)
        text(s, Inches(7.9), yy, Inches(0.55), rh,
             [{"segs": "›", "align": PP_ALIGN.CENTER}], color=acc, sz=22, bold=True,
             anchor=MSO_ANCHOR.MIDDLE)
    # right module: quick-ref
    m,_ = module(s, Inches(8.95), Inches(1.95), Inches(3.8), Inches(5.05),
                 title="DECK // QUICK REF")
    text(s, Inches(9.15), Inches(2.5), Inches(3.45), Inches(4.4),
         [
           {"segs":[("LEADER  ",{"color":MUTED,"font":MONO,"sz":10}),
                    ("Monkey D. Luffy (UG)",{"color":INK,"sz":12,"bold":True})],"space_after":7},
           {"segs":[("WIN-CON ",{"color":MUTED,"font":MONO,"sz":10}),
                    ("Boa Hancock + Navy HQ",{"color":GOLD,"sz":12,"bold":True})],"space_after":7},
           {"segs":[("IDENTITY",{"color":MUTED,"font":MONO,"sz":10})],"space_after":2},
           {"segs":"Aggro a valore sostenuto: vai largo, cicli, ri-erigi 2 DON!! ogni turno.",
            "color":INK_DIM,"sz":11,"space_after":8},
           {"segs":[("CORE    ",{"color":MUTED,"font":MONO,"sz":10})],"space_after":2},
           {"segs":"Luffy · Mr.1 · Mr.2 · Ivankov · Crocodile · Mr.3 · Buggy · Prisoner ×12",
            "color":INK_DIM,"sz":11,"space_after":8},
         ], line_spacing=1.18)
    led_bar(s, Inches(9.15), Inches(6.4), Inches(3.4), Inches(0.2), 5, 5, GREEN,
            peak_gold=True)
    text(s, Inches(9.15), Inches(6.62), Inches(3.4), Inches(0.3),
         [{"segs":"CONSISTENCY // MAX","track":120}], color=GREEN, sz=9,
         font=MONO, bold=True)
    footer_hud(s)


# ============================================================
#  DIVIDER (chapter break)
# ============================================================
def divider(num, title, sub, art_key, lead=False):
    s = slide(); base(s, vignette=False)
    acc = CHAP_ACCENT[num]
    # art a destra in scrim
    art_band(s, art_key, Inches(7.4), 0, Inches(5.93), SH, lead=lead, scrim=False)
    rect(s, Inches(7.4), 0, Inches(5.93), SH, BG_CARBON, alpha=42)
    rect(s, Inches(7.4), 0, Inches(2.6), SH, BG_DEEP, alpha=55)  # scrim sx dell'art
    rect(s, 0, 0, Inches(7.4), SH, BG_CARBON, alpha=30)
    numtxt = num if isinstance(num,str) else "%02d"%num
    # giant numeral EROE sul lato artwork (visibile, accento di capitolo)
    numcol = RGBColor(min(acc[0]+10,255), min(acc[1]+10,255), min(acc[2]+10,255))
    text(s, Inches(7.7), Inches(1.4), Inches(5.5), Inches(4.8), [numtxt],
         color=acc, sz=300, bold=True, font=DISP, align=PP_ALIGN.CENTER,
         anchor=MSO_ANCHOR.MIDDLE, wrap=False)
    rect(s, Inches(7.4), 0, Inches(5.93), SH, BG_DEEP, alpha=30)  # vela sul numeral
    # numeral inciso texture a sinistra
    giant(s, Inches(0.05), Inches(5.6), numtxt, sz=210, color=GRID_MAJOR,
          align=PP_ALIGN.LEFT)
    # rail verticale + label
    rect(s, Inches(0.7), Inches(2.5), Pt(2.4), Inches(2.4), acc)
    node_square(s, Inches(0.7)+Pt(1.2), Inches(2.5), size=Pt(12), color=acc)
    node_square(s, Inches(0.7)+Pt(1.2), Inches(4.9), size=Pt(12), color=acc)
    text(s, Inches(0.95), Inches(2.35), Inches(6), Inches(0.3),
         [{"segs": f"SECTOR {('%02d'%num) if not isinstance(num,str) else num}/06",
           "track": 220}], color=acc, sz=11, bold=True, font=MONO)
    text(s, Inches(0.92), Inches(2.75), Inches(6.4), Inches(1.4),
         [{"segs": title, "track": -10}], color=INK, sz=46, bold=True, font=DISP,
         line_spacing=0.95)
    text(s, Inches(0.95), Inches(4.25), Inches(6.0), Inches(0.6),
         [sub], color=INK_DIM, sz=14)
    gradient_band(s, SH-Inches(0.14), Inches(0.14))
    PNUM[0] += 1


# ============================================================
#  SLIDE 04 — CAP.1 IDENTITY · VALUE ENGINE
# ============================================================
def _engine_node(s, x, y, w, h, code, key, title, micro, accent):
    m = rect(s, x, y, w, h, SURFACE, line=accent, lw=1.1, shadow=True)
    rect(s, x, y, w, Inches(0.3), accent, alpha=22)
    text(s, x+Inches(0.12), y+Inches(0.04), w-Inches(0.2), Inches(0.28),
         [{"segs":[(code+"  ",{"font":MONO,"color":accent,"sz":9,"bold":True}),
                   (title,{"color":INK,"sz":12,"bold":True,"font":UISB})]}],
         anchor=MSO_ANCHOR.MIDDLE)
    # mini art chip
    pic_cover(s, A.card_art(key), x+Inches(0.12), y+Inches(0.38),
              Inches(0.92), Inches(0.52))
    corner_brackets(s, x+Inches(0.12), y+Inches(0.38), Inches(0.92), Inches(0.52),
                    accent, seg=Inches(0.1), lw=1.0)
    text(s, x+Inches(1.15), y+Inches(0.36), w-Inches(1.25), Inches(0.58),
         [micro], color=INK_DIM, sz=10, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.05)
    return m


def cap1_identity():
    s = slide(); base(s)
    kicker_h2(s, "SECTOR 01 // DECK IDENTITY",
              "VALUE ENGINE — NON GLASS CANNON", folio="DOC 01/06",
              sub="Il deck è un MOTORE a ciclo chiuso: ogni turno sviluppa, cicla, ri-erige e pressa.")
    giant(s, Inches(5.3), Inches(4.0), "∞", sz=260, color=GRID_MAJOR,
          align=PP_ALIGN.CENTER)

    # --- LOOP 4 nodi ---
    nw, nh = Inches(3.55), Inches(1.05)
    TLx, TLy = Inches(0.7), Inches(2.05)
    TRx, TRy = Inches(5.05), Inches(2.05)
    BLx, BLy = Inches(0.7), Inches(3.55)
    BRx, BRy = Inches(5.05), Inches(3.55)
    _engine_node(s, TLx, TLy, nw, nh, "01", "luffy1c", "DEVELOP",
                 "Gioca largo: 2-cost + Luffy 1c. Più corpi possibili.", BLUE)
    _engine_node(s, TRx, TRy, nw, nh, "02", "ivankov", "CYCLE",
                 "Quasi ogni carta si auto-pesca. Searcher + draw.", BLUE)
    _engine_node(s, BRx, BRy, nw, nh, "03", "crocodile", "REBUILD",
                 "Re-stand 2 DON!! · Crocodile bounce-replay.", GREEN)
    _engine_node(s, BLx, BLy, nw, nh, "04", "leader", "PRESS",
                 "Più swing nello stesso turno. Riduci mano + life.", GREEN)
    # connettori orari: 01->02 (blu, top), 02->03 (blu, dx giu), 03->04 (verde, bottom), 04->01 (verde, sx su)
    midy_t = TLy+nh/2; midy_b = BLy+nh/2
    arrow_flow(s, TLx+nw, midy_t, TRx, midy_t, BLUE, elbow="h")      # 01->02
    arrow_flow(s, TRx+nw/2, TRy+nh, TRx+nw/2, BRy, BLUE, elbow="v") # 02->03
    arrow_flow(s, BRx, midy_b, BLx+nw, midy_b, GREEN, elbow="h")     # 03->04 (verso sx)
    arrow_flow(s, BLx+nw/2, BLy, BLx+nw/2, TLy+nh, GREEN, elbow="v") # 04->01
    # engine core ovale centrale
    eng = rect(s, Inches(5.55), Inches(3.0), Inches(2.55), Inches(1.05),
               SURFACE_HI, line=CYAN, lw=1.3, shape=MSO_SHAPE.ROUNDED_RECTANGLE,
               radius=0.3, shadow=True)
    shape_text(eng, [{"segs":"VALUE ENGINE","sz":13,"bold":True,"color":CYAN,
                      "font":UISB,"track":80},
                     {"segs":"sustained, non one-shot","sz":9.5,"color":INK_DIM,
                      "font":MONO}])

    # --- colonna dx: 3 MODI ---
    m,_ = module(s, Inches(8.45), Inches(2.05), Inches(4.3), Inches(2.55),
                 title="OPERATING MODES // 1 ENGINE")
    modes = [("AGGRESSIVE","Pressione dal turno 1, swing larghi",BLUE,5),
             ("CONTROLLED","Freeza la board, gestisci il tempo",GREEN,3),
             ("FLEXIBLE","Tante tech card → adatti la build",GOLD,4)]
    yy = Inches(2.6)
    for nm,d,col,lv in modes:
        text(s, Inches(8.65), yy, Inches(2.1), Inches(0.3),
             [{"segs":nm,"track":60}], color=col, sz=12, bold=True, font=UISB)
        led_bar(s, Inches(10.85), yy+Inches(0.04), Inches(1.7), Inches(0.16), lv, 5, col)
        text(s, Inches(8.65), yy+Inches(0.28), Inches(4.0), Inches(0.3),
             [d], color=INK_DIM, sz=10.5)
        yy += Inches(0.62)

    # --- basamento CONSISTENCY ---
    bm = rect(s, Inches(0.7), Inches(4.85), Inches(7.4), Inches(1.55), SURFACE,
              line=GREEN, lw=1.1, shadow=True)
    rect(s, Inches(0.7), Inches(4.85), Inches(0.06), Inches(1.55), GREEN)
    text(s, Inches(0.9), Inches(4.96), Inches(7.0), Inches(0.3),
         [{"segs":"FOUNDATION // CONSISTENCY = FATTORE #1","track":140}],
         color=GREEN, sz=11, bold=True, font=UISB)
    text(s, Inches(0.9), Inches(5.32), Inches(7.05), Inches(1.0),
         [{"segs":[("Il deck gioca tante carte → ",{}),
                   ("altissima consistenza",{"bold":True,"color":INK}),
                   (". Quasi ogni carta si cicla da sola: se l'avversario insegue la board, "
                    "tu generi risorse e mantieni un aggro sostenuto, erodendo mano e life.",{})]}],
         color=INK_DIM, sz=12, line_spacing=1.15)
    # KPI a dx del basamento
    km,_ = module(s, Inches(8.45), Inches(4.85), Inches(4.3), Inches(1.55))
    data_big(s, Inches(8.45), Inches(4.95), Inches(2.1), "×12", "PRISONER", BLUE,
             vsz=46, sub="counter wall")
    data_big(s, Inches(10.6), Inches(4.95), Inches(2.1), "+2", "DON!! / TURN", GREEN,
             vsz=46, sub="re-stand leader")
    footer_hud(s, chap=1)


# ============================================================
#  SLIDE 06 — CAP.2 META OVERVIEW
# ============================================================
def cap2_meta_overview():
    s = slide(); base(s)
    kicker_h2(s, "SECTOR 02 // META ANALYSIS", "DA OP15 A OP16: LA META SI APRE",
              folio="DOC 02/06", accent=GREEN)
    # transition OP15 -> OP16
    o15,_ = module(s, Inches(0.55), Inches(1.85), Inches(3.15), Inches(1.5),
                   title="OP15 // META CHIUSA", accent=MUTED)
    text(s, Inches(0.75), Inches(2.35), Inches(2.8), Inches(0.9),
         [{"segs":"Dominata da 2 deck","sz":11,"color":INK_DIM,"space_after":4},
          {"segs":[("Purple Enel",{"color":INK,"bold":True}),
                   (" + ",{"color":MUTED}),("Nami",{"color":INK,"bold":True})],"sz":14}],
         anchor=MSO_ANCHOR.MIDDLE)
    chev(s, Inches(3.95), Inches(2.6), GOLD, "right", size=Inches(0.2))
    o16,_ = module(s, Inches(4.25), Inches(1.85), Inches(8.5), Inches(1.5),
                   title="OP16 // META « JUNGLE »", accent=GOLD, hi=True)
    text(s, Inches(4.5), Inches(2.32), Inches(8.1), Inches(1.0),
         [{"segs":[("Tanti contender, più deck tier-1, molto matchup-dependent. ",
                    {"color":INK}),
                   ("Catalizzatore: arrivo di BY Teach",{"color":GOLD,"bold":True}),
                   (" che counterа sia Enel sia Nami.",{"color":INK})]}],
         color=INK_DIM, sz=12.5, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.15)

    # roster 6 contender
    text(s, Inches(0.55), Inches(3.55), Inches(8), Inches(0.3),
         [{"segs":"FIELD // 6 CONTENDER PRINCIPALI","track":160}], color=CYAN,
         sz=11, bold=True, font=MONO)
    six = [("teach","TEACH","BY · popolare",DANGER,False),
           ("nami","NAMI","UY · sempre",BLUE,False),
           ("rosinante","ROSINANTE","PY · top",RGBColor(0x9B,0x5C,0xE0),False),
           ("luffy","LUFFY","UG · NOI",GREEN,True),
           ("enel","ENEL","P · in calo",GOLD,False),
           ("yamato","YAMATO","B · tricky",INK_DIM,False)]
    n=6; x0=Inches(0.55); tot=Inches(12.2); gap=Inches(0.16)
    bw=int((tot-gap*(n-1))/n)
    for i,(k,nm,d,col,us) in enumerate(six):
        x=x0+(bw+gap)*i
        h=Inches(2.5); y=Inches(3.9)
        m=rect(s, x, y, bw, h, SURFACE_HI if us else SURFACE,
               line=col, lw=2.0 if us else 1.0, shadow=True)
        pic_cover(s, A.lead_art(k), x, y, bw, Inches(1.35))
        scrim_band(s, x, y+Inches(0.7), bw, Inches(0.65), "up")
        if us:
            corner_brackets(s, x, y, bw, h, GREEN, seg=Inches(0.16), lw=1.5)
            chip(s, x+Inches(0.1), y+Inches(0.1), bw-Inches(0.2), Inches(0.3),
                 "◆ US", GREEN, BG_CARBON, sz=9.5, track=80)
        text(s, x+Inches(0.1), y+Inches(1.4), bw-Inches(0.2), Inches(0.34),
             [{"segs":nm,"track":40,"align":PP_ALIGN.CENTER}],
             color=(GOLD if us else INK), sz=13, bold=True, font=UISB)
        text(s, x+Inches(0.06), y+Inches(1.75), bw-Inches(0.12), Inches(0.6),
             [{"segs":d,"align":PP_ALIGN.CENTER}], color=INK_DIM, sz=10,
             anchor=MSO_ANCHOR.TOP)
        led_dot(s, x+bw/2-Inches(0.065), y+Inches(2.22),
                state=("green" if us else "muted"))
    # nota jungle
    text(s, Inches(0.55), Inches(6.55), Inches(12.2), Inches(0.4),
         [{"segs":[("REGOLA JUNGLE:  ",{"color":GOLD,"bold":True,"font":MONO,"sz":11}),
                   ("ogni deck ha in media 4 buoni matchup e 2 cattivi.  UG Luffy fa da "
                    "gatekeeper a molti deck tier-2.",{"color":INK_DIM,"sz":12})]}],
         line_spacing=1.0)
    footer_hud(s, chap=2)


# ============================================================
#  SLIDE 07 — CAP.2 TIER LIST
# ============================================================
def cap2_tierlist():
    s = slide(); base(s)
    kicker_h2(s, "SECTOR 02 // TIER ASSESSMENT", "TIER LIST OP16 (PRO PLAYER JP)",
              folio="DOC 02/06", accent=GREEN)
    # screenshot in frame
    nx,ny,nw,nh = pic_contain(s, A.img("ch2_tierlist_op16.png"),
                              Inches(0.55), Inches(1.85), Inches(7.6), Inches(5.0))
    corner_brackets(s, nx-Pt(3), ny-Pt(3), nw+Pt(6), nh+Pt(6), CYAN,
                    seg=Inches(0.24), lw=1.6)
    scanline(s, nx, ny+nh-Inches(0.28), nw, "INTEL // TIER-LIST.JP", color=CYAN)
    # take autore
    m,_ = module(s, Inches(8.45), Inches(1.85), Inches(4.3), Inches(3.3),
                 title="AUTHOR // TAKE", accent=GREEN)
    text(s, Inches(8.65), Inches(2.35), Inches(3.95), Inches(2.75),
         [{"segs":"La tier list riflette solo in parte la realtà.","sz":12,
           "bold":True,"color":INK,"space_after":7},
          {"segs":[("▸ ",{"color":GREEN,"bold":True}),
                   ("Meta più jungle-like: molti contender, più tier-1.",{})],
           "sz":11,"color":INK_DIM,"space_after":6},
          {"segs":[("▸ ",{"color":GREEN,"bold":True}),
                   ("Enel scende: lo counterano Teach e Nami.",{})],
           "sz":11,"color":INK_DIM,"space_after":6},
          {"segs":[("▸ ",{"color":GREEN,"bold":True}),
                   ("Yamato = tier-1 più debole, ma molto tricky.",{})],
           "sz":11,"color":INK_DIM,"space_after":6},
          {"segs":[("▸ ",{"color":GREEN,"bold":True}),
                   ("Si rimodellerà in base a popolarità e risultati.",{})],
           "sz":11,"color":INK_DIM}],
         line_spacing=1.15)
    term_log(s, Inches(8.45), Inches(5.4), Inches(4.3),
             [("MONITOR: cosa è popolare / cosa cala", INK_DIM),
              ("BE READY: fine-tune della lista", GREEN)],
             title="SYS // RACCOMANDAZIONE", accent=GREEN)
    footer_hud(s, chap=2)


# ============================================================
#  SLIDE 08 — CAP.2 WIN-RATE MATRIX (regina)
# ============================================================
def cap2_winrate():
    s = slide(); base(s)
    kicker_h2(s, "SECTOR 02 // WIN-RATE MATRIX",
              "MATRICE 6×6 — SIM, SETTIMANA 1", folio="DOC 02/06", accent=GREEN,
              sub="% vittoria del leader di RIGA vs colonna · ogni cella: barra BLU = andando 1°, barra VERDE = andando 2°.")
    giant(s, Inches(10.2), Inches(1.5), "6×6", sz=150, color=GRID_MAJOR,
          align=PP_ALIGN.LEFT)
    rows = ["TEACH","YAMATO","LUFFY","ROSIN.","NAMI","ENEL"]
    cols = ["TEA","YAM","LUF","ROS","NAM","ENE"]
    data = [
        [(25,75),(24,47),(25,41),(42,53),(60,77),(70,66)],
        [(53,76),(44,56),(48,62),(44,42),(27,31),(47,38)],
        [(59,75),(39,52),(43,57),(56,44),(34,29),(38,32)],
        [(47,58),(59,56),(56,44),(59,41),(49,36),(42,31)],
        [(23,40),(69,73),(71,67),(64,51),(47,53),(76,61)],
        [(33,30),(62,53),(68,62),(69,58),(39,24),(58,42)],
    ]
    LUF = 2  # indice riga/colonna Luffy
    ox, oy = Inches(0.55), Inches(2.0)
    hdr_w, hdr_h = Inches(1.45), Inches(0.5)
    cw, ch = Inches(1.52), Inches(0.66)
    # angolo
    text(s, ox, oy, hdr_w, hdr_h, [{"segs":"ROW vs ▸","align":PP_ALIGN.CENTER}],
         color=MUTED, sz=10, bold=True, font=MONO, anchor=MSO_ANCHOR.MIDDLE)
    # header colonne
    for c,cn in enumerate(cols):
        x = ox+hdr_w+cw*c
        isL = (c==LUF)
        b=rect(s, x, oy, cw-Pt(2), hdr_h, SURFACE_HI if isL else SURFACE,
               line=GRID_MAJOR, lw=0.6)
        shape_text(b, cn, color=(GOLD if isL else INK), sz=11, bold=True,
                   font=MONO, track=40)
    # righe
    for r,rn in enumerate(rows):
        y = oy+hdr_h+ch*r
        isLr = (r==LUF)
        b=rect(s, ox, y, hdr_w, ch-Pt(2),
               SURFACE_HI if isLr else SURFACE, line=GRID_MAJOR, lw=0.6)
        shape_text(b, [{"segs":rn,"align":PP_ALIGN.LEFT}],
                   color=(GOLD if isLr else INK), sz=11, bold=True, font=MONO,
                   align=PP_ALIGN.LEFT, track=20)
        if isLr:
            rect(s, ox, y, Inches(0.05), ch-Pt(2), GREEN)
        for c,(w1,w2) in enumerate(data[r]):
            x = ox+hdr_w+cw*c
            cell = rect(s, x, y, cw-Pt(2), ch-Pt(2),
                        SURFACE if (r+c)%2==0 else BG_CARBON,
                        line=GRID_MAJOR, lw=0.5)
            if r==c:  # mirror diagonale
                rect(s, x, y, cw-Pt(2), ch-Pt(2), GRID_MAJOR, alpha=40)
                shape_text(cell, "MIRROR", color=MUTED, sz=8.5, font=MONO, track=40)
                continue
            # tint verdetto sulla riga Luffy
            if isLr:
                avg=(w1+w2)/2
                vc = GREEN if avg>=52 else (DANGER if avg<45 else BLUE)
                rect(s, x, y, cw-Pt(2), ch-Pt(2), vc, alpha=14)
            # 2 micro led-bar
            bw=cw-Inches(0.5)
            led_bar(s, x+Inches(0.05), y+Inches(0.12), bw, Inches(0.1),
                    round(w1/10), 10, BLUE)
            led_bar(s, x+Inches(0.05), y+Inches(0.40), bw, Inches(0.1),
                    round(w2/10), 10, GREEN)
            text(s, x+bw+Inches(0.06), y, Inches(0.42), ch-Pt(2),
                 [{"segs":str(w1),"color":BLUE,"sz":9,"font":MONO,"bold":True,"space_after":1},
                  {"segs":str(w2),"color":GREEN,"sz":9,"font":MONO,"bold":True}],
                 anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.LEFT)
    # identity row brackets sulla riga Luffy intera
    ry = oy+hdr_h+ch*LUF
    corner_brackets(s, ox, ry, hdr_w+cw*6-Pt(2), ch-Pt(2), GREEN,
                    seg=Inches(0.16), lw=1.6)

    # legenda + KPI a destra
    m,_ = module(s, Inches(10.3), Inches(3.0), Inches(2.45), Inches(2.0),
                 title="LEGEND", accent=CYAN)
    text(s, Inches(10.5), Inches(3.5), Inches(2.2), Inches(1.4),
         [{"segs":[("▰ ",{"color":BLUE}),("GO 1ST",{"color":INK_DIM})],"sz":10.5,"font":MONO,"space_after":5},
          {"segs":[("▰ ",{"color":GREEN}),("GO 2ND",{"color":INK_DIM})],"sz":10.5,"font":MONO,"space_after":8},
          {"segs":[("■ ",{"color":GREEN}),("GOOD ≥52%",{"color":INK_DIM})],"sz":10,"font":MONO,"space_after":4},
          {"segs":[("■ ",{"color":DANGER}),("BAD <45%",{"color":INK_DIM})],"sz":10,"font":MONO}])
    term_log(s, Inches(0.55), Inches(6.5), Inches(12),
             [("LUFFY ROW: ottimo vs Teach (59/75) · vs Nami/Enel « un po' off » → più vincibili dei numeri", INK_DIM)],
             title=None, accent=GREEN)
    text(s, Inches(0.55), Inches(6.82), Inches(12), Inches(0.3),
         [{"segs":"⚠ DATI SIM SETT.1 · LISTE NON AFFINATE · DA PRENDERE CON LE PINZE","track":120}],
         color=MUTED, sz=8.5, font=MONO)
    footer_hud(s, chap=2)


# ============================================================
#  SLIDE 10 — CAP.3 DECKLIST GRID (carte reali)
# ============================================================
def cap3_decklist():
    s = slide(); base(s)
    kicker_h2(s, "SECTOR 03 // DECK SCHEMATIC", "LA LISTA GENERICA — 50 CARTE",
              folio="DOC 03/06",
              sub="Start di meta: lista il più GENERICA possibile. Focus su VALUE GENERATION + board che pressa.")
    # 12 carte staple, 6 per riga x 2
    grid = [("luffy1c","4"),("mr1","4"),("mr2","4"),("ivankov","4"),
            ("crocodile","4"),("mr3","4"),("buggy","4"),("prisoner","12"),
            ("hancock","3"),("samurai","4"),("navyhq","2"),("gumgumrain","1")]
    cols=6; cw=Inches(1.9); gx=Inches(0.18); x0=Inches(0.6)
    cardw=Inches(1.55); cardh=cardw*1.396
    for i,(k,cnt) in enumerate(grid):
        r=i//cols; c=i%cols
        x=x0+(cw+gx)*c; y=Inches(1.95)+(cardh+Inches(0.55))*r
        gold = (k=="hancock")
        nx,ny,nw,nh = pic_contain(s, A.card_full(k), x, y, cardw, cardh)
        corner_brackets(s, nx-Pt(2), ny-Pt(2), nw+Pt(4), nh+Pt(4),
                        gold=gold, seg=Inches(0.13), lw=1.2)
        # badge copie
        bd=rect(s, nx+nw-Inches(0.5), ny-Inches(0.12), Inches(0.6), Inches(0.46),
                GOLD if gold else SURFACE_HI, line=(GOLD if gold else CYAN), lw=1.2,
                shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.3, shadow=True)
        shape_text(bd, "×"+cnt, color=(BG_CARBON if gold else INK), sz=13, bold=True,
                   font=DISP)
        # scanline codice
        text(s, nx, ny+nh+Pt(2), nw, Inches(0.2),
             [{"segs":A.code(k),"align":PP_ALIGN.CENTER,"track":40}],
             color=(GOLD if gold else MUTED), sz=8, font=MONO, bold=True)
    footer_hud(s, chap=3)


# ============================================================
#  CARD SPOTLIGHT (scheda dettaglio riutilizzabile)
# ============================================================
def card_slide(kicker, keys, title, role, points, verdict, copies, combo=None):
    """keys: lista di 1-2 chiavi carta da mostrare affiancate."""
    s = slide(); base(s)
    gold = ("hancock" in keys)
    kicker_h2(s, kicker, title, folio="DOC 03/06",
              accent=(GOLD if gold else CYAN))
    # area immagini sinistra
    if len(keys)==1:
        card_spotlight(s, keys[0], Inches(0.7), Inches(1.95), Inches(3.3),
                       Inches(4.6), gold=gold, code=A.code(keys[0]))
    else:
        card_spotlight(s, keys[0], Inches(0.6), Inches(2.05), Inches(2.45),
                       Inches(3.4), code=A.code(keys[0]))
        card_spotlight(s, keys[1], Inches(2.35), Inches(2.55), Inches(2.45),
                       Inches(3.4), code=A.code(keys[1]))
    # giant copies
    giant(s, Inches(0.7), Inches(6.7), copies, sz=64,
          color=(GOLD if gold else GRID_MAJOR), align=PP_ALIGN.LEFT)
    # right column
    rx = Inches(5.15); rw = Inches(7.6)
    # role + cost chips
    chip(s, rx, Inches(1.95), Inches(2.6), Inches(0.42),
         role, SURFACE, (GOLD if gold else BLUE), line=(GOLD if gold else BLUE),
         sz=11.5, track=80)
    for i,k in enumerate(keys):
        cost_chip(s, rx+Inches(2.75)+Inches(1.7)*i, Inches(1.99), k)
    # points module
    m,top = module(s, rx, Inches(2.55), rw, Inches(2.5) if combo else Inches(3.2),
                   title="PERCHÉ È FORTE", accent=(GOLD if gold else CYAN))
    yy = Inches(3.05)
    for pt in points:
        led_dot(s, rx+Inches(0.22), yy+Inches(0.05), state=("gold" if gold else "blue"))
        text(s, rx+Inches(0.48), yy-Inches(0.04), rw-Inches(0.7), Inches(0.5),
             [pt], color=INK_DIM, sz=12.5, line_spacing=1.08)
        yy += Inches(0.52)
    # verdict
    vy = Inches(5.2) if combo else Inches(5.9)
    if not combo:
        vm = rect(s, rx, vy, rw, Inches(0.95), SURFACE_HI,
                  line=(GOLD if gold else GREEN), lw=1.2, shadow=True)
        rect(s, rx, vy, Inches(0.06), Inches(0.95), GOLD if gold else GREEN)
        text(s, rx+Inches(0.2), vy+Inches(0.08), rw-Inches(0.4), Inches(0.3),
             [{"segs":"VERDETTO","track":160}], color=(GOLD if gold else GREEN),
             sz=10, bold=True, font=MONO)
        text(s, rx+Inches(0.2), vy+Inches(0.34), rw-Inches(0.4), Inches(0.55),
             [verdict], color=INK, sz=12.5, bold=True, line_spacing=1.05)
    if combo:
        cm = rect(s, rx, Inches(5.2), rw, Inches(1.45), SURFACE,
                  line=GOLD, lw=1.3, shadow=True)
        rect(s, rx, Inches(5.2), Inches(2.0), Inches(1.45), GOLD, alpha=12)
        text(s, rx+Inches(0.2), Inches(5.28), Inches(1.9), Inches(1.3),
             [{"segs":"COMBO","sz":17,"bold":True,"color":GOLD,"font":DISP},
              {"segs":"LINE","sz":17,"bold":True,"color":GOLD,"font":DISP}],
             anchor=MSO_ANCHOR.MIDDLE)
        text(s, rx+Inches(2.1), Inches(5.3), rw-Inches(2.3), Inches(1.3),
             [combo], color=INK_DIM, sz=11.5, anchor=MSO_ANCHOR.MIDDLE,
             line_spacing=1.12)
        # verdict mini sotto i punti
        text(s, rx, Inches(5.02), rw, Inches(0.2),
             [{"segs":[("VERDETTO  ",{"color":GREEN,"font":MONO,"sz":9,"bold":True,"track":120}),
                       (verdict,{"color":INK,"sz":11,"bold":True})]}],
             line_spacing=1.0)
    footer_hud(s, chap=3)


# ============================================================
#  SLIDE 19 — CAP.3 TECH POOL (semaforo)
# ============================================================
def _tech_row(s, x, y, w, key_or_none, nm, desc):
    h=Inches(0.86)
    rect(s, x, y, w, h, BG_CARBON, alpha=40)
    if key_or_none:
        pic_cover(s, A.card_art(key_or_none), x+Inches(0.08), y+Inches(0.1),
                  Inches(0.95), Inches(0.66))
        corner_brackets(s, x+Inches(0.08), y+Inches(0.1), Inches(0.95),
                        Inches(0.66), CYAN, seg=Inches(0.1), lw=1.0)
        tx=x+Inches(1.18)
    else:
        tx=x+Inches(0.18)
    text(s, tx, y+Inches(0.08), w-(tx-x)-Inches(0.1), Inches(0.3),
         [{"segs":nm,"sz":11.5,"bold":True,"color":INK}])
    text(s, tx, y+Inches(0.38), w-(tx-x)-Inches(0.1), Inches(0.46),
         [desc], color=INK_DIM, sz=9.5, line_spacing=1.0)


def cap3_techcards():
    s = slide(); base(s)
    kicker_h2(s, "SECTOR 03 // TECH POOL", "IL POOL DI TECH CARD — VERDETTI",
              folio="DOC 03/06",
              sub="Punto di forza del deck: enorme varietà di tool card da inserire in base alla meta attesa.")
    cols=[("IN / IN TEST",GREEN,"green",[
            ("mr3_2c","Mr.3 Galdino 2k counter","Abilita lethal, ricicla DON!! · in test ×2"),
            ("iknow","I Know You're Strong...","Freeza 2 corpi · ottima in mirror & vs Roronoa"),
            ("mr3_2c","Mr.3 Galdino (2c)","Toolbox: nome diverso per il re-stand")]),
          ("SITUAZIONALI",GOLD,"gold",[
            (None,"DMG (Dead Man's Game)","1-of vs deck difensivi (nega Blocker)"),
            ("brook","Brook","Event-like: rest + re-stand · combo Hancock"),
            (None,"Hound Blaze","Fortissima solo in mirror"),
            (None,"Inazuma","Se serve Blocker piccolo ed economico")]),
          ("OUT / SCARTATE",DANGER,"danger",[
            ("jinbe","Jinbe","2k counter ma = -1, non serve"),
            (None,"Ivankov 3c / Mr.2 4c / Luffy 4c","Troppo costose / clunky"),
            (None,"Miss Olive · Crocodile 6c","Lente, corpo debole"),
            (None,"Gum-Gum Champion Rifle","Troppo random · Electrical Luna < I Know")])]
    x0=Inches(0.55); cw=Inches(4.0); gx=Inches(0.13)
    for i,(title,col,state,items) in enumerate(cols):
        x=x0+(cw+gx)*i
        m=rect(s, x, Inches(2.0), cw, Inches(4.65), SURFACE, line=col, lw=1.2,
               shadow=True)
        rect(s, x, Inches(2.0), cw, Inches(0.5), col, alpha=20)
        rect(s, x, Inches(2.0), Inches(0.06), Inches(4.65), col)
        led_dot(s, x+Inches(0.18), Inches(2.17), state=state)
        text(s, x+Inches(0.42), Inches(2.06), cw-Inches(0.5), Inches(0.4),
             [{"segs":title,"track":120}], color=col, sz=13, bold=True, font=UISB,
             anchor=MSO_ANCHOR.MIDDLE)
        yy=Inches(2.62)
        for k,nm,d in items:
            _tech_row(s, x+Inches(0.12), yy, cw-Inches(0.24), k, nm, d)
            yy+=Inches(0.96)
    footer_hud(s, chap=3)


# ============================================================
#  SLIDE 21 — CAP.4 PLAYSTYLE
# ============================================================
def cap4_playstyle():
    s = slide(); base(s)
    kicker_h2(s, "SECTOR 04 // PLAYSTYLE", "LO STILE DI GIOCO CORRETTO",
              folio="DOC 04/06", accent=GREEN)
    art_band(s, "leader", Inches(9.5), Inches(1.95), Inches(3.25), Inches(2.0),
             code="OP16-022", scrim=True)
    corner_brackets(s, Inches(9.5), Inches(1.95), Inches(3.25), Inches(2.0), CYAN,
                    seg=Inches(0.18), lw=1.3)
    # MITO vs REALTA
    e=rect(s, Inches(0.55), Inches(1.95), Inches(4.25), Inches(1.5), SURFACE,
           line=DANGER, lw=1.3, shadow=True)
    rect(s, Inches(0.55), Inches(1.95), Inches(4.25), Inches(0.42), DANGER, alpha=18)
    text(s, Inches(0.72), Inches(2.0), Inches(4.0), Inches(0.34),
         [{"segs":"✗  L'ERRORE COMUNE","track":80}], color=DANGER, sz=12, bold=True,
         font=UISB, anchor=MSO_ANCHOR.MIDDLE)
    text(s, Inches(0.72), Inches(2.45), Inches(3.95), Inches(0.9),
         [{"segs":"« È un glass cannon → aggro immediato a tutti i costi »"}],
         color=INK_DIM, sz=12.5, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.1)
    chev(s, Inches(4.97), Inches(2.7), GOLD, "right", size=Inches(0.22))
    rr=rect(s, Inches(5.3), Inches(1.95), Inches(4.0), Inches(1.5), SURFACE_HI,
            line=GREEN, lw=1.3, shadow=True)
    rect(s, Inches(5.3), Inches(1.95), Inches(4.0), Inches(0.42), GREEN, alpha=18)
    text(s, Inches(5.47), Inches(2.0), Inches(3.7), Inches(0.34),
         [{"segs":"✓  LA REALTÀ","track":80}], color=GREEN, sz=12, bold=True,
         font=UISB, anchor=MSO_ANCHOR.MIDDLE)
    text(s, Inches(5.47), Inches(2.45), Inches(3.7), Inches(0.9),
         [{"segs":"Sustained VALUE ENGINE: costruisci la board, poi swing tutti insieme."}],
         color=INK, sz=12.5, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.1)
    # 4 principi
    text(s, Inches(0.55), Inches(3.6), Inches(8), Inches(0.3),
         [{"segs":"CORE PRINCIPLES // 4","track":160}], color=CYAN, sz=11,
         bold=True, font=MONO)
    pr=[("VAI LARGO","Più corpi, restano stanti, poi attacchi insieme. Pochi board wipe.",GREEN,5),
        ("SWING MULTIPLI","Conta il NUMERO di attacchi, non la dimensione.",BLUE,5),
        ("RESTA HEALTHY","Counter gli attacchi da 1 carta. Baita l'avversario nella board.",GOLD,4),
        ("RI-COSTRUISCI","Esponi la board sapendo di rifarla (Crocodile, Buggy, draw).",BLUE,4)]
    n=4;x0=Inches(0.55);tot=Inches(12.2);gap=Inches(0.18)
    bw=int((tot-gap*(n-1))/n)
    for i,(t,d,col,lv) in enumerate(pr):
        x=x0+(bw+gap)*i
        m=rect(s, x, Inches(3.95), bw, Inches(1.95), SURFACE, line=col, lw=1.1,
               shadow=True)
        rect(s, x, Inches(3.95), bw, Inches(0.06), col)
        text(s, x+Inches(0.15), Inches(4.05), Inches(0.7), Inches(0.5),
             [str(i+1)], color=col, sz=30, bold=True, font=DISP)
        text(s, x+Inches(0.85), Inches(4.12), bw-Inches(0.95), Inches(0.4),
             [{"segs":t,"track":40}], color=INK, sz=13, bold=True, font=UISB,
             anchor=MSO_ANCHOR.MIDDLE)
        text(s, x+Inches(0.15), Inches(4.65), bw-Inches(0.3), Inches(0.9),
             [d], color=INK_DIM, sz=11, line_spacing=1.1)
        led_bar(s, x+Inches(0.15), Inches(5.55), bw-Inches(0.3), Inches(0.14), lv, 5, col)
    # life note
    lb=rect(s, Inches(0.55), Inches(6.1), Inches(12.2), Inches(0.78), SURFACE,
            line=GOLD, lw=1.1)
    rect(s, Inches(0.55), Inches(6.1), Inches(0.06), Inches(0.78), GOLD)
    text(s, Inches(0.78), Inches(6.18), Inches(11.9), Inches(0.65),
         [{"segs":[("LIFE MANAGEMENT  ",{"color":GOLD,"font":MONO,"sz":10,"bold":True,"track":100}),
                   ("Stare a 4 life è pericoloso. Il deck pesca tanto → counterare gli attacchi early e "
                    "preservare i life baita l'avversario nella board e ti fa guadagnare tempo.",
                    {"color":INK_DIM,"sz":11.5})]}],
         anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.08)
    footer_hud(s, chap=4)


# ============================================================
#  SLIDE 22 — CAP.4 MULLIGAN
# ============================================================
def cap4_mulligan():
    s = slide(); base(s)
    kicker_h2(s, "SECTOR 04 // MULLIGAN", "COSA CERCARE IN APERTURA",
              folio="DOC 04/06", accent=GREEN,
              sub="Obiettivo: poter sviluppare più corpi nei primi 2 turni.")
    # target cards row (mini art chips)
    targets=[("luffy1c","Luffy 1c","cerca quasi tutto",BLUE),
             ("mr1","Mr.1","opener + counter",BLUE),
             ("mr2","Mr.2","no counter → gioca",BLUE),
             ("samurai","Samurai","combo su 2 corpi",GREEN),
             ("crocodile","Crocodile/Ivankov","4c per andare 2°",GOLD)]
    n=5;x0=Inches(0.55);tot=Inches(12.2);gap=Inches(0.15)
    bw=int((tot-gap*(n-1))/n)
    for i,(k,nm,d,col) in enumerate(targets):
        x=x0+(bw+gap)*i
        m=rect(s, x, Inches(1.95), bw, Inches(1.5), SURFACE, line=col, lw=1.0, shadow=True)
        pic_cover(s, A.card_art(k), x+Inches(0.08), Inches(2.05), bw-Inches(0.16), Inches(0.72))
        corner_brackets(s, x+Inches(0.08), Inches(2.05), bw-Inches(0.16),
                        Inches(0.72), col, seg=Inches(0.1), lw=1.0)
        text(s, x+Inches(0.1), Inches(2.83), bw-Inches(0.2), Inches(0.3),
             [{"segs":nm,"align":PP_ALIGN.CENTER}], color=INK, sz=11, bold=True,
             font=UISB)
        text(s, x+Inches(0.06), Inches(3.12), bw-Inches(0.12), Inches(0.3),
             [{"segs":d,"align":PP_ALIGN.CENTER}], color=INK_DIM, sz=9)
    # curve pari/dispari
    ev,_=module(s, Inches(0.55), Inches(3.6), Inches(6.0), Inches(1.3),
                title="CURVA PARI // ANDANDO 2° (PREFERITA)", accent=GREEN)
    text(s, Inches(0.78), Inches(4.05), Inches(5.6), Inches(0.8),
         ["2× corpi 1c/2c → re-stand DON!! → Samurai. Poi Crocodile/Ivankov per sviluppare e cercare follow-up."],
         color=INK_DIM, sz=11.5, line_spacing=1.12)
    od,_=module(s, Inches(6.75), Inches(3.6), Inches(6.0), Inches(1.3),
                title="CURVA DISPARI // ANDANDO 1°", accent=BLUE)
    text(s, Inches(6.98), Inches(4.05), Inches(5.6), Inches(0.8),
         ["Buggy 5c + Prisoner sul turno da 5. Crocodile forte 1° con un 2c. Boa Hancock sul turno da 7."],
         color=INK_DIM, sz=11.5, line_spacing=1.12)
    # 4 esempi keep/mull
    text(s, Inches(0.55), Inches(5.02), Inches(8), Inches(0.3),
         [{"segs":"HAND EXAMPLES // 4","track":160}], color=CYAN, sz=11, bold=True,
         font=MONO)
    ex=[("ES.1 · 1°","MULLIGAN",DANGER,"Nessun gioco turni 1-3 (1 e 3 DON!! vuoti)."),
        ("ES.2 · 1°","KEEP",GREEN,"Luffy 1c → Mr.1 → Mr.2 + Samurai: insane."),
        ("ES.3 · 2°","KEEP",GREEN,"Come es.2 + Ivankov: massima flessibilità."),
        ("ES.4 · 2°","KEEP",GREEN,"Mr.2 + Crocodile + Buggy follow-up: solida.")]
    bw=int((tot-gap*(n-1))/n)  # riuso 5? no, 4
    n4=4; bw4=int((tot-Inches(0.18)*(n4-1))/n4)
    for i,(t,verd,col,d) in enumerate(ex):
        x=x0+(bw4+Inches(0.18))*i
        m=rect(s, x, Inches(5.4), bw4, Inches(1.45), SURFACE, line=col, lw=1.1,
               shadow=True)
        text(s, x+Inches(0.15), Inches(5.5), bw4-Inches(0.3), Inches(0.3),
             [{"segs":t,"track":60}], color=MUTED, sz=10, bold=True, font=MONO)
        chip(s, x+Inches(0.15), Inches(5.82), bw4-Inches(0.3), Inches(0.4),
             verd, col, BG_CARBON if col!=DANGER else INK, sz=12, track=60)
        text(s, x+Inches(0.15), Inches(6.3), bw4-Inches(0.3), Inches(0.5),
             [d], color=INK_DIM, sz=10, line_spacing=1.05)
    footer_hud(s, chap=4)


# ============================================================
#  SLIDE 23 — CAP.4 GAMEPLAN (phase rail)
# ============================================================
def cap4_gameplan():
    s = slide(); base(s)
    kicker_h2(s, "SECTOR 04 // GAME PLAN", "EARLY → MID → LATE → KILL",
              folio="DOC 04/06", accent=GREEN)
    giant(s, Inches(10.5), Inches(1.4), "RUN", sz=120, color=GRID_MAJOR,
          align=PP_ALIGN.LEFT)
    phases=[("01","EARLY","Costruisci la board",GREEN,"green",
             ["Sviluppa 2c + Luffy 1c, vai largo","2c → re-stand → Samurai (+2)",
              "Buggy + Prisoner a 5 = 2 corpi 6k"],"ch4_board_early.png",3),
            ("02","MID","Imposta il lethal",BLUE,"blue",
             ["Non smettere di sviluppare","Re-stand 2 DON!! → push + sviluppi",
              "Boa Hancock freeza il corpo chiave"],"ch4_board_mid.png",5),
            ("03","LATE","Chiudi la partita",GOLD,"gold",
             ["Swing sul Leader power avversario","Re-stand Event: 4-5 nomi diversi",
              "Navy HQ → ri-erigi e vai lethal"],"ch4_board_endgame.png",7)]
    x0=Inches(0.55); cw=Inches(3.75); gx=Inches(0.42)
    rail_y=Inches(2.15)
    for i,(num,t,sub,col,st,pts,img,lv) in enumerate(phases):
        x=x0+(cw+gx)*i
        # node header
        node_square(s, x+Inches(0.2), rail_y, size=Pt(12), color=col)
        text(s, x+Inches(0.42), rail_y-Inches(0.16), Inches(2), Inches(0.32),
             [{"segs":f"PHASE {num}","track":160}], color=col, sz=10, bold=True,
             font=MONO, anchor=MSO_ANCHOR.MIDDLE)
        m=rect(s, x, rail_y+Inches(0.25), cw, Inches(4.4), SURFACE, line=col,
               lw=1.2, shadow=True)
        rect(s, x, rail_y+Inches(0.25), cw, Inches(0.7), col, alpha=16)
        text(s, x+Inches(0.18), rail_y+Inches(0.3), cw-Inches(0.3), Inches(0.4),
             [{"segs":t,"track":40}], color=INK, sz=20, bold=True, font=DISP)
        text(s, x+Inches(0.18), rail_y+Inches(0.72), cw-Inches(0.3), Inches(0.3),
             [sub], color=col, sz=10.5, font=MONO)
        # board screenshot
        nx,ny,nw,nh=pic_contain(s, A.img(img), x+Inches(0.15),
                                rail_y+Inches(1.1), cw-Inches(0.3), Inches(1.35))
        corner_brackets(s, nx, ny, nw, nh, col, seg=Inches(0.14), lw=1.1)
        yy=rail_y+Inches(2.65)
        for pt in pts:
            led_dot(s, x+Inches(0.2), yy+Inches(0.03), state=st, size=Inches(0.1))
            text(s, x+Inches(0.4), yy-Inches(0.04), cw-Inches(0.55), Inches(0.45),
                 [pt], color=INK_DIM, sz=10.5, line_spacing=1.05)
            yy+=Inches(0.46)
        # status led bar (don progression)
        text(s, x+Inches(0.18), rail_y+Inches(4.05), Inches(1.4), Inches(0.25),
             [{"segs":"DON!!","track":80}], color=MUTED, sz=8.5, font=MONO)
        led_bar(s, x+Inches(0.95), rail_y+Inches(4.07), cw-Inches(1.15),
                Inches(0.14), lv, 10, col, peak_gold=(st=="gold"))
        # connettore tra fasi
        if i<2:
            ax=x+cw; ay=rail_y+Inches(2.4)
            rect(s, ax, ay-Pt(1.5), gx, Pt(3), col)
            chev(s, ax+gx, ay, col, "right", size=Inches(0.16))
    # KILL node
    kx=x0+(cw+gx)*3-Inches(0.05)
    footer_hud(s, chap=4)


# ============================================================
#  SLIDE 24 — COMBO LINE (Hancock + Mr.3 + Samurai)
# ============================================================
def combo_line():
    s = slide(); base(s)
    kicker_h2(s, "SECTOR 04 // SIGNATURE LINE",
              "COMBO: HANCOCK + MR.3 + SAMURAI", folio="DOC 04/06", accent=GOLD,
              sub="Un turno-chiave che spende DON!! a step per PESCARE 4 e FREEZARE 2.")
    giant(s, Inches(11.1), Inches(1.5), "9", sz=150, color=GRID_MAJOR,
          align=PP_ALIGN.LEFT)
    steps=[("01","hancock","BOA HANCOCK","Gioca 7c · freeza 1 · corpo 9k",GOLD,True,"OP16-032","blue"),
           ("02","mr3","MR.3 GALDINO","Re-stand 2 DON!! · pesca 2 · freeza",BLUE,False,"OP16-056","green"),
           ("03","samurai","YOU CAN BE MY SAMURAI","Resti Hancock+Mr.3 · re-stand · pesca 2",GREEN,False,"OP01-055","blue")]
    x0=Inches(0.55); cw=Inches(3.35); gx=Inches(0.5); ny=Inches(2.3)
    railc=ny+Inches(1.7)
    for i,(num,k,nm,d,col,gold,code,inconn) in enumerate(steps):
        x=x0+(cw+gx)*i
        m=rect(s, x, ny, cw, Inches(3.3), SURFACE, line=(GOLD if gold else col),
               lw=1.4 if gold else 1.1, shadow=True)
        text(s, x+Inches(0.15), ny+Inches(0.08), Inches(1.2), Inches(0.45),
             [num], color=(GOLD if gold else col), sz=26, bold=True, font=DISP)
        text(s, x+Inches(1.0), ny+Inches(0.12), cw-Inches(1.1), Inches(0.45),
             [{"segs":nm,"track":20}], color=INK, sz=12.5, bold=True, font=UISB,
             anchor=MSO_ANCHOR.MIDDLE)
        # card art
        pic_cover(s, A.card_art(k), x+Inches(0.15), ny+Inches(0.62),
                  cw-Inches(0.3), Inches(1.55))
        corner_brackets(s, x+Inches(0.15), ny+Inches(0.62), cw-Inches(0.3),
                        Inches(1.55), gold=gold, seg=Inches(0.14), lw=1.2)
        scanline(s, x+Inches(0.15), ny+Inches(1.93), cw-Inches(0.3), code,
                 color=(GOLD if gold else CYAN))
        text(s, x+Inches(0.18), ny+Inches(2.35), cw-Inches(0.36), Inches(0.85),
             [d], color=INK_DIM, sz=11, line_spacing=1.12)
        # connettore
        if i<2:
            ax=x+cw; cc=GREEN if steps[i+1][7]=="green" else BLUE
            rect(s, ax, railc-Pt(1.5), gx, Pt(3), cc)
            node_square(s, ax, railc, size=Pt(8), color=cc)
            chev(s, ax+gx, railc, cc, "right", size=Inches(0.15))
    # OUTPUT module
    ox=x0+(cw+gx)*3-Inches(0.1)
    om=rect(s, Inches(0.55), Inches(5.85), Inches(8.6), Inches(1.05), SURFACE_HI,
            line=GOLD, lw=1.3, shadow=True)
    rect(s, Inches(0.55), Inches(5.85), Inches(0.06), Inches(1.05), GOLD)
    text(s, Inches(0.78), Inches(5.93), Inches(3), Inches(0.3),
         [{"segs":"RESOLVE // OUTPUT","track":140}], color=GOLD, sz=10, bold=True,
         font=MONO)
    data_big(s, Inches(0.7), Inches(6.05), Inches(1.8), "+4", "PESCA", BLUE, vsz=40)
    data_big(s, Inches(2.5), Inches(6.05), Inches(1.8), "2", "FREEZE", DANGER, vsz=40)
    text(s, Inches(4.5), Inches(5.95), Inches(4.5), Inches(0.9),
         [{"segs":"Pesca 4 carte e congela 2 personaggi nemici in un solo turno. Enorme value, e Hancock resta come win-con."}],
         color=INK_DIM, sz=11.5, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.12)
    # DON budget
    dm,_=module(s, Inches(9.35), Inches(5.85), Inches(3.4), Inches(1.05),
                title="DON!! BUDGET", accent=BLUE)
    led_bar(s, Inches(9.55), Inches(6.45), Inches(3.0), Inches(0.22), 7, 9, BLUE)
    text(s, Inches(9.55), Inches(6.66), Inches(3.0), Inches(0.25),
         [{"segs":"9 → 2 reali (re-stand paga il resto)","track":60}], color=INK_DIM,
         sz=9, font=MONO)
    footer_hud(s, chap=4)


# ============================================================
#  LIFE HEARTS (mini-dashboard vite)
# ============================================================
def life_hearts(s, x, y, keep_n, total=4, danger_below=2):
    sz=Inches(0.34); gap=Inches(0.1)
    for i in range(total):
        lifeval=total-i  # mostra 4..1
        cx=x+(sz+gap)*i
        on = lifeval<=keep_n
        col = DANGER if (on and lifeval<=danger_below) else (INK if on else MUTED)
        fill = (col if on else BG_DEEP)
        b=rect(s, cx, y, sz, sz, fill if on else BG_DEEP,
               line=(None if on else GRID_MAJOR), lw=0.75,
               shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.3)
        shape_text(b, str(lifeval), color=(BG_CARBON if on else MUTED), sz=13,
                   bold=True, font=DISP)


# ============================================================
#  SLIDE 26 — CAP.5 MATCHUP OVERVIEW
# ============================================================
def cap5_overview():
    s = slide(); base(s)
    kicker_h2(s, "SECTOR 05 // MATCHUP INDEX", "QUADRO DEI 6 MATCHUP",
              folio="DOC 05/06")
    data=[("luffy","MIRROR (BG LUFFY)","EVEN","2°","Board control + life. Si decide nel combat."),
          ("teach","BY TEACH","FAVOR.","2°","« Siediti e vinci ». Vai largo, non trascinare."),
          ("rosinante","PY ROSINANTE","EVEN","1°","Non vincere la board war: svuota la mano."),
          ("nami","UY NAMI","EVEN","1°","Corpi 6k che swingano 7k. Hancock chiave."),
          ("yamato","B YAMATO","EVEN","2°","Entrambi al Leader. Resta healthy, clock."),
          ("enel","P ENEL","EVEN","1°","Corpi sticky, aggro grezzo. No Hancock.")]
    cols=3; cw=Inches(4.0); gx=Inches(0.13); x0=Inches(0.55)
    cy=[Inches(1.95), Inches(4.5)]
    for idx,(k,nm,diff,start,desc) in enumerate(data):
        c=idx%3; r=idx//3
        x=x0+(cw+gx)*c; y=cy[r]
        dc = GREEN if diff=="FAVOR." else GOLD
        m=rect(s, x, y, cw, Inches(2.4), SURFACE, line=GRID_MAJOR, lw=0.75, shadow=True)
        # leader art chip
        nx,ny,nw,nh=pic_contain(s, A.lead_full(k), x+Inches(0.12), y+Inches(0.14),
                                Inches(1.15), Inches(1.6))
        corner_brackets(s, nx-Pt(2), ny-Pt(2), nw+Pt(4), nh+Pt(4), CYAN,
                        seg=Inches(0.12), lw=1.1)
        text(s, x+Inches(1.45), y+Inches(0.18), cw-Inches(1.6), Inches(0.4),
             [{"segs":nm,"track":20}], color=INK, sz=13, bold=True, font=UISB,
             anchor=MSO_ANCHOR.MIDDLE)
        chip(s, x+Inches(1.45), y+Inches(0.68), Inches(1.3), Inches(0.36),
             diff, dc, BG_CARBON, sz=10.5, track=60)
        chip(s, x+Inches(2.85), y+Inches(0.68), Inches(0.95), Inches(0.36),
             "ST "+start, SURFACE_HI, INK, line=CYAN, sz=10, track=40, font=MONO)
        text(s, x+Inches(1.45), y+Inches(1.18), cw-Inches(1.6), Inches(1.1),
             [desc], color=INK_DIM, sz=10.5, line_spacing=1.1)
        led_dot(s, x+Inches(1.45), y+Inches(2.05),
                state=("green" if diff=="FAVOR." else "gold"))
    footer_hud(s, chap=5)


# ============================================================
#  MATCHUP TEMPLATE
# ============================================================
def matchup_slide(key, name, diff, start, overview, early, mid, late, tips, life,
                  keep_n, life_note=None, ring=CYAN):
    s = slide(); base(s)
    diff_disp = "FAVOREVOLE" if "Easy" in diff or "favor" in diff.lower() else "EVEN"
    dc = GREEN if diff_disp=="FAVOREVOLE" else GOLD
    kicker_h2(s, "SECTOR 05 // MATCHUP", name, folio="DOC 05/06", accent=dc)
    # COLONNA INTEL (sx)
    nx,ny,nw,nh=pic_contain(s, A.lead_full(key), Inches(0.55), Inches(1.85),
                            Inches(2.05), Inches(2.85))
    corner_brackets(s, nx-Pt(3), ny-Pt(3), nw+Pt(6), nh+Pt(6), ring,
                    seg=Inches(0.2), lw=1.6)
    scanline(s, nx, ny+nh-Inches(0.28), nw, A.LEADERS[key][0], color=ring)
    chip(s, Inches(0.55), Inches(4.85), Inches(2.05), Inches(0.42),
         diff_disp, dc, BG_CARBON, sz=12, track=80)
    chip(s, Inches(0.55), Inches(5.32), Inches(2.05), Inches(0.42),
         "START: "+start, SURFACE_HI, INK, line=CYAN, sz=11, track=60, font=MONO)
    om,_=module(s, Inches(0.55), Inches(5.8), Inches(2.05), Inches(1.05),
                title="OVERVIEW", accent=dc)
    text(s, Inches(0.72), Inches(6.2), Inches(1.8), Inches(0.62),
         [overview], color=INK_DIM, sz=9.5, line_spacing=1.05)
    # COLONNA GAME PLAN (rail verticale 3 fasi)
    phases=[("EARLY",GREEN,"green",early),("MID",BLUE,"blue",mid),("LATE",GOLD,"gold",late)]
    px=Inches(2.85); pw=Inches(6.25); yy=Inches(1.85); ph=Inches(1.18)
    for i,(t,col,st,txt_) in enumerate(phases):
        m=rect(s, px, yy, pw, ph, SURFACE, line=col, lw=1.0, shadow=True)
        tab=rect(s, px, yy, Inches(1.1), ph, col, alpha=16)
        rect(s, px, yy, Inches(0.06), ph, col)
        text(s, px+Inches(0.15), yy, Inches(1.0), ph,
             [{"segs":t,"track":40}], color=col, sz=14, bold=True, font=DISP,
             anchor=MSO_ANCHOR.MIDDLE)
        led_dot(s, px+Inches(0.15), yy+ph-Inches(0.24), state=st, size=Inches(0.1))
        text(s, px+Inches(1.25), yy+Inches(0.04), pw-Inches(1.4), ph-Inches(0.08),
             [txt_], color=INK_DIM, sz=11, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.08)
        if i<2:
            cc=phases[i+1][1]
            node_square(s, px+Inches(0.55), yy+ph, size=Pt(7), color=cc)
            rect(s, px+Inches(0.55)-Pt(1.2), yy+ph, Pt(2.4), Inches(0.08), cc)
        yy+=ph+Inches(0.08)
    # KEY TIPS (dx)
    km,_=module(s, Inches(9.35), Inches(1.85), Inches(3.4), Inches(3.96),
                title="KEY TIPS", accent=GOLD)
    lines=[]
    for tp in tips:
        lines.append({"segs":[("▸ ",{"color":GOLD,"bold":True}),(tp,{})],
                      "sz":10.5,"color":INK_DIM,"space_after":5})
    text(s, Inches(9.55), Inches(2.35), Inches(3.05), Inches(3.4), lines,
         line_spacing=1.1)
    # LIFE MANAGEMENT banda inferiore
    ly=Inches(5.95); lh=Inches(0.92)
    lb=rect(s, Inches(2.85), ly, Inches(9.9), lh, SURFACE, line=DANGER, lw=1.3,
            shadow=True)
    rect(s, Inches(2.85), ly, Inches(2.0), lh, DANGER, alpha=14)
    rect(s, Inches(2.85), ly, Inches(0.06), lh, DANGER)
    text(s, Inches(3.0), ly+Inches(0.08), Inches(1.85), Inches(0.34),
         [{"segs":"♥ LIFE MGMT","track":60}], color=DANGER, sz=12, bold=True,
         font=UISB)
    text(s, Inches(3.0), ly+Inches(0.42), Inches(1.85), Inches(0.26),
         [{"segs":"keep ~"+str(keep_n)+" life","track":40}], color=MUTED, sz=9,
         font=MONO)
    life_hearts(s, Inches(3.0), ly+Inches(0.62), keep_n)
    # life bullets
    text(s, Inches(4.95), ly+Inches(0.07), Inches(7.7), lh-Inches(0.14),
         [{"segs":[("• ",{"color":DANGER,"bold":True}),(it,{})],"sz":10.5,
           "color":INK_DIM,"space_after":2} for it in life],
         anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.06)
    footer_hud(s, chap=5)
